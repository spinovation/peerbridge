import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set to 16:9 widescreen layout
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # HSL-mapped Premium Color Palette
    BG_COLOR = RGBColor(11, 15, 26)         # Deep Navy (#0b0f1a)
    SURFACE_COLOR = RGBColor(22, 28, 45)    # Slate Card Surface (#161c2d)
    TEXT_MAIN = RGBColor(255, 255, 255)     # Crisp White (#ffffff)
    TEXT_MUTED = RGBColor(163, 163, 163)    # Muted Grey (#a3a3a3)
    CYBER_BLUE = RGBColor(0, 242, 254)      # Primary Cyber Blue (#00f2fe)
    EMERALD_GREEN = RGBColor(16, 185, 129)  # Vetted Emerald Green (#10b981)
    ROSE_RED = RGBColor(244, 63, 94)        # Warning Red (#f43f5e)
    
    blank_slide_layout = prs.slide_layouts[6] # Blank slide layout

    def set_dark_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # ---------------------------------------------------------
    # SLIDE 1: Title Slide (Cover)
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    set_dark_bg(slide)
    
    # Accent glowing rect on top border
    glowing_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.12))
    glowing_bar.fill.solid()
    glowing_bar.fill.fore_color.rgb = CYBER_BLUE
    glowing_bar.line.fill.background()

    # Title Text Frame
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "PEER BRIDGE"
    p.font.name = "Inter"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = CYBER_BLUE
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "PRIVATE DEBT & EQUITY ECOSYSTEM"
    p2.font.name = "Inter"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_MAIN
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "Full Technical Architecture Masterclass (Phase 1, Phase 2 & Phase 3)"
    p3.font.name = "Inter"
    p3.font.size = Pt(18)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(30)

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.5), Inches(11.33), Inches(0.5))
    ftf = footer_box.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "PROPRIETARY & CONFIDENTIAL  •  PEER BRIDGE NETWORKS INC."
    fp.font.name = "Inter"
    fp.font.size = Pt(10)
    fp.font.bold = True
    fp.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 2: Executive Summary
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    set_dark_bg(slide)
    
    # Title
    t_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.33), Inches(0.8))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "EXECUTIVE SUMMARY"
    p.font.name = "Inter"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CYBER_BLUE
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Building decentralized commercial debt syndicates & retail equity marketplaces"
    p_sub.font.name = "Inter"
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.space_before = Pt(5)

    # Core concept list
    c_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(6.0), Inches(4.5))
    ctf = c_box.text_frame
    ctf.word_wrap = True
    
    bullets = [
        ("The Capital Underwriting Paradigm Shift", "Bypassing FICO scores which act as lagging indicators and penalize asset-rich but credit-thin founders. Integrating real-time paycheck structures and transactional analysis."),
        ("Next-Generation Financial Vetting", "Leveraging ADP & Paychex Payroll APIs for gross-to-net salary verification alongside Plaid transactional category parsing to calculate true discretionary cash flow."),
        ("End-to-End Automation Pipeline", "Automating fractional Escrow Splits, Auto-Invest Sweeps, IRS Form 1099-INT interest tax compilation, and SEC Reg CF Form C filing sheets."),
        ("Phase 3 Autonomous AI Brokerage", "Launching agent nodes (CapitalAgent, FounderAgent, AuditAgent) to negotiate rates, evaluate prospectuses, and auto-verify KYC credentials on the blockchain ledger.")
    ]
    
    for i, (title, desc) in enumerate(bullets):
        bp = ctf.add_paragraph() if i > 0 else ctf.paragraphs[0]
        bp.text = f"•  {title}"
        bp.font.name = "Inter"
        bp.font.size = Pt(16)
        bp.font.bold = True
        bp.font.color.rgb = CYBER_BLUE
        if i > 0: bp.space_before = Pt(15)
        
        dp = ctf.add_paragraph()
        dp.text = desc
        dp.font.name = "Inter"
        dp.font.size = Pt(11)
        dp.font.color.rgb = TEXT_MUTED
        dp.space_before = Pt(2)
        dp.level = 0

    # Decorative Panel: Technical Core Pillars (Right Side)
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(1.8), Inches(4.8), Inches(4.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = SURFACE_COLOR
    panel.line.color.rgb = RGBColor(38, 48, 77)
    
    ptf = panel.text_frame
    ptf.word_wrap = True
    ptf.margin_left = Inches(0.4)
    ptf.margin_right = Inches(0.4)
    ptf.margin_top = Inches(0.4)
    
    pp = ptf.paragraphs[0]
    pp.text = "TECHNICAL CORE PILLARS"
    pp.font.name = "Inter"
    pp.font.size = Pt(14)
    pp.font.bold = True
    pp.font.color.rgb = EMERALD_GREEN
    
    pp2 = ptf.add_paragraph()
    pp2.text = "\n1. DECENTRALIZED DATA SWEEPS\nReal-time secure credentials handshake simulated directly in the client sandbox.\n\n2. AUTOMATED COMPLIANCE\nForm 1099-INT interactive preview compilers and one-click W-9 verification vaults.\n\n3. AI DEPLOYED AGENTS\nAutonomous negotiation routines generating cryptographic SHA-256 agreement signatures.\n\n4. PRODUCTION SPEED\nBuilt under Next.js Turbopack with 3.2s compile profiles and zero runtime alerts."
    pp2.font.name = "Inter"
    pp2.font.size = Pt(11)
    pp2.font.color.rgb = TEXT_MAIN
    pp2.space_before = Pt(5)

    # ---------------------------------------------------------
    # SLIDE 3: Phase 1 - Foundational Infrastructure
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    set_dark_bg(slide)
    
    t_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.33), Inches(0.8))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PHASE 1: Foundational Infrastructure & Design"
    p.font.name = "Inter"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CYBER_BLUE
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Premium design systems and core peer-to-peer transaction components"
    p_sub.font.name = "Inter"
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.space_before = Pt(5)

    # Three key modules columns
    modules = [
        ("Design System (HSL)", "Cyberpunk Glassmorphism", "Harmonized Dark HSL variables mapping back-drops, border-effects, glows, and responsive typography.\n\n• Background: 224 25% 4%\n• Card: 224 25% 6%\n• Border: 220 20% 12%\n• Glows: backdrop-filter blur"),
        ("Fundraising Listings", "P2P Commercial Notes", "Fully interactive campaign lists with dynamic progress bars, ledger fractional purchases, and escrow wallet settlements.\n\n• P2P Commercial Notes\n• Y-Combinator SAFE equity\n• live escrow allocations"),
        ("Network Directory", "Symmetrical P2P Handshake", "Global search and directory with instant biometrics handshake matching, active connection states, and LinkedIn DMs.\n\n• Symmetrical DM channels\n• Credentials status checks\n• In-app sandbox requests")
    ]
    
    for idx, (m_title, m_sub, m_desc) in enumerate(modules):
        left_pos = Inches(1.0 + idx * 3.9)
        col_panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(1.8), Inches(3.53), Inches(4.5))
        col_panel.fill.solid()
        col_panel.fill.fore_color.rgb = SURFACE_COLOR
        col_panel.line.color.rgb = RGBColor(38, 48, 77)
        
        col_tf = col_panel.text_frame
        col_tf.word_wrap = True
        col_tf.margin_left = Inches(0.25)
        col_tf.margin_right = Inches(0.25)
        col_tf.margin_top = Inches(0.3)
        
        mp1 = col_tf.paragraphs[0]
        mp1.text = m_title
        mp1.font.name = "Inter"
        mp1.font.size = Pt(16)
        mp1.font.bold = True
        mp1.font.color.rgb = CYBER_BLUE
        
        mp2 = col_tf.add_paragraph()
        mp2.text = m_sub
        mp2.font.name = "Inter"
        mp2.font.size = Pt(11)
        mp2.font.bold = True
        mp2.font.color.rgb = EMERALD_GREEN
        mp2.space_before = Pt(2)
        
        mp3 = col_tf.add_paragraph()
        mp3.text = "\n" + m_desc
        mp3.font.name = "Inter"
        mp3.font.size = Pt(10.5)
        mp3.font.color.rgb = TEXT_MUTED
        mp3.space_before = Pt(5)

    # ---------------------------------------------------------
    # SLIDE 4: Phase 2 - Advanced PRI Risk Underwriting
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    set_dark_bg(slide)
    
    t_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.33), Inches(0.8))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PHASE 2: Proprietary Risk Index (PRI) Credit Engine"
    p.font.name = "Inter"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CYBER_BLUE
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Four-layer underwriting scoring matrix with real-time API integrations"
    p_sub.font.name = "Inter"
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.space_before = Pt(5)

    # Text box for Layers description
    l_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.8), Inches(4.5))
    ltf = l_box.text_frame
    ltf.word_wrap = True
    
    layers = [
        ("Layer 0: Identity & Fraud", "Sanctions/PEP screens, device fingerprints, and synthetic identity risk anomalies."),
        ("Layer 1: Bureau BCS", "Bureau credit tier (300-850), tradelines count, utilization ratios, and inquiry audit."),
        ("Layer 2: BRS (Behavioral Risk Score)", "Connects payroll (ADP/Paychex) and Plaid bank feeds to isolate gross-to-net paychecks, tax withholdings, mandatory obligations, and discretionary spends."),
        ("Layer 3: Public Records Check", "Tax liens, legal judgments, historical bankruptcies, and years since discharge.")
    ]
    
    for i, (l_title, l_desc) in enumerate(layers):
        lp = ltf.add_paragraph() if i > 0 else ltf.paragraphs[0]
        lp.text = f"•  {l_title}"
        lp.font.name = "Inter"
        lp.font.size = Pt(14)
        lp.font.bold = True
        lp.font.color.rgb = CYBER_BLUE
        if i > 0: lp.space_before = Pt(12)
        
        ldp = ltf.add_paragraph()
        ldp.text = l_desc
        ldp.font.name = "Inter"
        ldp.font.size = Pt(10.5)
        ldp.font.color.rgb = TEXT_MUTED
        ldp.space_before = Pt(2)

    # Right side: Bypass explanation box
    bypass_panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(1.8), Inches(5.0), Inches(4.5))
    bypass_panel.fill.solid()
    bypass_panel.fill.fore_color.rgb = SURFACE_COLOR
    bypass_panel.line.color.rgb = RGBColor(38, 48, 77)
    
    bptf = bypass_panel.text_frame
    bptf.word_wrap = True
    bptf.margin_left = Inches(0.3)
    bptf.margin_right = Inches(0.3)
    bptf.margin_top = Inches(0.3)
    
    bp1 = bptf.paragraphs[0]
    bp1.text = "⚡ THE BUREAU-BYPASS PARADIGM"
    bp1.font.name = "Inter"
    bp1.font.size = Pt(15)
    bp1.font.bold = True
    bp1.font.color.rgb = EMERALD_GREEN
    
    bp2 = bptf.add_paragraph()
    bp2.text = "\nCredit bureaus are lagging databases. High-earning, cash-flow stable founders often have thin credit files or lack debt history, leading to low FICO ratings.\n\nOur modern credit engine introduces a dynamic override: if ADP/Paychex and Plaid accounts are successfully connected, traditional Bureau FICO weight drops to 10%, while cash flow metrics take 90% weight.\n\nBRS evaluates:\n•  Gross-to-net W-2 paycheck stability\n•  Pre-tax deductions (401k/HSA)\n•  Discretionary shopping burn velocity\n•  Debt-to-Disposable-Income (DDI) margins"
    bp2.font.name = "Inter"
    bp2.font.size = Pt(11)
    bp2.font.color.rgb = TEXT_MAIN
    bp2.space_before = Pt(5)

    # ---------------------------------------------------------
    # SLIDE 5: Case Study - Underwriting Vetting Contrast
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    set_dark_bg(slide)
    
    t_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.33), Inches(0.8))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "UNDERWRITING CONTRAST: High Earner vs Responsible Saver"
    p.font.name = "Inter"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CYBER_BLUE
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Dynamic Bureau-Bypass results demonstrating true cash flow risk"
    p_sub.font.name = "Inter"
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.space_before = Pt(5)

    # Card 1: Devon Vance (Decline)
    card1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.8), Inches(5.4), Inches(4.5))
    card1.fill.solid()
    card1.fill.fore_color.rgb = SURFACE_COLOR
    card1.line.color.rgb = ROSE_RED
    card1.line.width = Pt(2)
    
    c1_tf = card1.text_frame
    c1_tf.word_wrap = True
    c1_tf.margin_left = Inches(0.3)
    c1_tf.margin_top = Inches(0.3)
    
    cp1 = c1_tf.paragraphs[0]
    cp1.text = "DEVON VANCE (AURORA SYSTEMS)"
    cp1.font.name = "Inter"
    cp1.font.size = Pt(16)
    cp1.font.bold = True
    cp1.font.color.rgb = ROSE_RED
    
    cp1_sub = c1_tf.add_paragraph()
    cp1_sub.text = "High Income / High Lifestyle Burn"
    cp1_sub.font.name = "Inter"
    cp1_sub.font.size = Pt(10.5)
    cp1_sub.font.bold = True
    cp1_sub.font.color.rgb = TEXT_MUTED
    cp1_sub.space_before = Pt(2)
    
    cp1_body = c1_tf.add_paragraph()
    cp1_body.text = "\n•  Traditional FICO: 610 (Thin file / Fair)\n•  ADP Verified Annual Gross: $500,000\n•  California Tax Deductions: $18,000/mo\n•  Monthly Net Take-Home: $21,500/mo\n•  Plaid Mandatory Debts: $8,500/mo\n•  Plaid Discretionary Burn: $11,500/mo (Luxury)\n•  Revolving Cards: Maxed Out (85% Util)\n•  Monthly Savings Rate: 6.9% ($1,500 surplus)\n\nUNDERWRITING RESULT: DECLINED (PRI 520 - P5)\nSevere default hazard due to zero cushion reserves."
    cp1_body.font.name = "Inter"
    cp1_body.font.size = Pt(11)
    cp1_body.font.color.rgb = TEXT_MAIN
    cp1_body.space_before = Pt(5)

    # Card 2: Elena Rostova (Approve)
    card2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.4), Inches(4.5))
    card2.fill.solid()
    card2.fill.fore_color.rgb = SURFACE_COLOR
    card2.line.color.rgb = EMERALD_GREEN
    card2.line.width = Pt(2)
    
    c2_tf = card2.text_frame
    c2_tf.word_wrap = True
    c2_tf.margin_left = Inches(0.3)
    c2_tf.margin_top = Inches(0.3)
    
    cp2 = c2_tf.paragraphs[0]
    cp2.text = "ELENA ROSTOVA (NEUROWEB AI)"
    cp2.font.name = "Inter"
    cp2.font.size = Pt(16)
    cp2.font.bold = True
    cp2.font.color.rgb = EMERALD_GREEN
    
    cp2_sub = c2_tf.add_paragraph()
    cp2_sub.text = "Moderate Income / High Cash Cushion"
    cp2_sub.font.name = "Inter"
    cp2_sub.font.size = Pt(10.5)
    cp2_sub.font.bold = True
    cp2_sub.font.color.rgb = TEXT_MUTED
    cp2_sub.space_before = Pt(2)
    
    cp2_body = c2_tf.add_paragraph()
    cp2_body.text = "\n•  Traditional FICO: 620 (Thin file / Fair)\n•  ADP Verified Annual Gross: $160,000\n•  Tax Deductions: Optimized at $2,600/mo\n•  Monthly Net Take-Home: $9,200/mo\n•  Plaid Mandatory Debts: $1,800/mo\n•  Plaid Discretionary Burn: $1,100/mo (Disciplined)\n•  Revolving Cards: Zero Debt (15% Util)\n•  Monthly Savings Rate: 68.4% ($6,300 surplus)\n\nUNDERWRITING RESULT: APPROVED (PRI 780 - P2)\nFICO overridden. High capital limit issued due to massive cash flow buffer and strong financial discipline."
    cp2_body.font.name = "Inter"
    cp2_body.font.size = Pt(11)
    cp2_body.font.color.rgb = TEXT_MAIN
    cp2_body.space_before = Pt(5)

    # ---------------------------------------------------------
    # SLIDE 6: Phase 2 - Pro SaaS & AP Automation
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    set_dark_bg(slide)
    
    t_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.33), Inches(0.8))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PHASE 2: Premium SaaS & Automation Suites"
    p.font.name = "Inter"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CYBER_BLUE
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Automated tax systems, accounts payable, and crowd SPV prep panels"
    p_sub.font.name = "Inter"
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.space_before = Pt(5)

    # Columns description for SaaS & AP
    col1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.8), Inches(5.4), Inches(4.5))
    col1.fill.solid()
    col1.fill.fore_color.rgb = SURFACE_COLOR
    col1.line.color.rgb = RGBColor(38, 48, 77)
    
    col1_tf = col1.text_frame
    col1_tf.word_wrap = True
    col1_tf.margin_left = Inches(0.3)
    col1_tf.margin_top = Inches(0.3)
    
    c1p1 = col1_tf.paragraphs[0]
    c1p1.text = "LENDER PRO SYSTEM ($29/mo)"
    c1p1.font.name = "Inter"
    c1p1.font.size = Pt(16)
    c1p1.font.bold = True
    c1p1.font.color.rgb = CYBER_BLUE
    
    c1p2 = col1_tf.add_paragraph()
    c1p2.text = "\n•  Auto-Invest matching engine sweeps: Monitors launching commercial notes, matching risk bounds & target yields in real time.\n\n•  Escrow splitting: Executes direct ledger allocations (fractional notes split dynamically).\n\n•  Form 1099-INT Compiler: Automatically compiles interest income earned. Clicking generates a high-fidelity interactive IRS Form 1099-INT modal."
    c1p2.font.name = "Inter"
    c1p2.font.size = Pt(11)
    c1p2.font.color.rgb = TEXT_MAIN
    c1p2.space_before = Pt(5)

    col2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.4), Inches(4.5))
    col2.fill.solid()
    col2.fill.fore_color.rgb = SURFACE_COLOR
    col2.line.color.rgb = RGBColor(38, 48, 77)
    
    col2_tf = col2.text_frame
    col2_tf.word_wrap = True
    col2_tf.margin_left = Inches(0.3)
    col2_tf.margin_top = Inches(0.3)
    
    c2p1 = col2_tf.paragraphs[0]
    c2p1.text = "FOUNDER PRO AUTOMATION ($49/mo)"
    c2p1.font.name = "Inter"
    c2p1.font.size = Pt(16)
    c2p1.font.bold = True
    c2p1.font.color.rgb = EMERALD_GREEN
    
    c2p2 = col2_tf.add_paragraph()
    c2p2.text = "\n•  Accounts Payable AP Console: Outstanding invoice tracking, cash runway forecasting, and auto-bill adjustments.\n\n•  Chore Delegation: Delegates legal/accounting review chores to verified affiliates with escrow-locked incentives.\n\n•  SEC Form C Compliance Prep: Pre-populates W-2 disclosures, cap tables, and past round details into SEC template sheets."
    c2p2.font.name = "Inter"
    c2p2.font.size = Pt(11)
    c2p2.font.color.rgb = TEXT_MAIN
    c2p2.space_before = Pt(5)

    # ---------------------------------------------------------
    # SLIDE 7: Phase 3 - Autonomous AI Agents
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    set_dark_bg(slide)
    
    t_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.33), Inches(0.8))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PHASE 3: Autonomous AI Agent Brokerage"
    p.font.name = "Inter"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CYBER_BLUE
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Multi-agent autonomous brokerage, compliance auditing, and rate negotiation"
    p_sub.font.name = "Inter"
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.space_before = Pt(5)

    # Column 1: Core AI Agents
    c_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.8), Inches(4.5))
    ctf = c_box.text_frame
    ctf.word_wrap = True
    
    agents = [
        ("CapitalAgent (Investor AI)", "Parses startup offering prospectuses, evaluates ARR curves, checks BRS cash flows, and triggers fraction escrow bids on syndicates."),
        ("FounderAgent (Entrepreneur AI)", "Optimizes commercial note interest terms based on modern savings rate, updates cap tables, and manages accounts payable schedules."),
        ("AuditAgent (Affiliate Compliance AI)", "Performs passport liveness validation, accredited checks, and W-9 audits, claiming transaction commission fees.")
    ]
    
    for i, (a_title, a_desc) in enumerate(agents):
        ap = ctf.add_paragraph() if i > 0 else ctf.paragraphs[0]
        ap.text = f"•  {a_title}"
        ap.font.name = "Inter"
        ap.font.size = Pt(14)
        ap.font.bold = True
        ap.font.color.rgb = CYBER_BLUE
        if i > 0: ap.space_before = Pt(15)
        
        adp = ctf.add_paragraph()
        adp.text = a_desc
        adp.font.name = "Inter"
        adp.font.size = Pt(10.5)
        adp.font.color.rgb = TEXT_MUTED
        adp.space_before = Pt(2)

    # Column 2: Negotiation Simulator Panel
    neg_panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(1.8), Inches(5.0), Inches(4.5))
    neg_panel.fill.solid()
    neg_panel.fill.fore_color.rgb = SURFACE_COLOR
    neg_panel.line.color.rgb = RGBColor(38, 48, 77)
    
    ntf = neg_panel.text_frame
    ntf.word_wrap = True
    ntf.margin_left = Inches(0.3)
    ntf.margin_right = Inches(0.3)
    ntf.margin_top = Inches(0.3)
    
    np1 = ntf.paragraphs[0]
    np1.text = "🤝 INTERACTIVE RATE NEGOTIATOR"
    np1.font.name = "Inter"
    np1.font.size = Pt(15)
    np1.font.bold = True
    np1.font.color.rgb = EMERALD_GREEN
    
    np2 = ntf.add_paragraph()
    np2.text = "\nWithin AIAgentHub.js, users can trigger a live negotiation simulator:\n\n1. Pitch: FounderAgent requests a commercial note at 6.5% APR based on their optimized BRS cash flow.\n\n2. Counter: CapitalAgent analyzes FICO metrics, counter-pitching 9.5% APR to mitigate thin credit.\n\n3. Settlement: Agents run mathematical calculations, converging on an optimum 7.8% APR.\n\n4. Immutable Vault: Generates a Promissory Note and prints a secure SHA-256 cryptographic signature."
    np2.font.name = "Inter"
    np2.font.size = Pt(11)
    np2.font.color.rgb = TEXT_MAIN
    np2.space_before = Pt(5)

    # ---------------------------------------------------------
    # SLIDE 8: Tech Stack & Verification
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)
    set_dark_bg(slide)
    
    t_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.33), Inches(0.8))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Tech Stack, Build Metrics & Vetting"
    p.font.name = "Inter"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CYBER_BLUE
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Next.js production compiler benchmarks and code base layout"
    p_sub.font.name = "Inter"
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.space_before = Pt(5)

    # Column 1: Tech Stack
    ts_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.8), Inches(4.5))
    tstf = ts_box.text_frame
    tstf.word_wrap = True
    
    tech_stacks = [
        ("Premium Web Framework", "Next.js 16 (Turbopack) & React client-side hooks"),
        ("Global State Engine", "usePeerBridge.js custom telemetry controller"),
        ("Core Styling System", "Custom CSS Variables mapped in HSL dark variables"),
        ("Relational Database Path", "Firebase Cloud Firestore (11 Collections)")
    ]
    
    for i, (t_title, t_desc) in enumerate(tech_stacks):
        tp = tstf.add_paragraph() if i > 0 else tstf.paragraphs[0]
        tp.text = f"•  {t_title}"
        tp.font.name = "Inter"
        tp.font.size = Pt(14)
        tp.font.bold = True
        tp.font.color.rgb = CYBER_BLUE
        if i > 0: tp.space_before = Pt(15)
        
        tdp = tstf.add_paragraph()
        tdp.text = t_desc
        tdp.font.name = "Inter"
        tdp.font.size = Pt(10.5)
        tdp.font.color.rgb = TEXT_MUTED
        tdp.space_before = Pt(2)

    # Column 2: Build Metrics
    metrics_panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(1.8), Inches(5.0), Inches(4.5))
    metrics_panel.fill.solid()
    metrics_panel.fill.fore_color.rgb = SURFACE_COLOR
    metrics_panel.line.color.rgb = RGBColor(38, 48, 77)
    
    mtf = metrics_panel.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0.3)
    mtf.margin_right = Inches(0.3)
    mtf.margin_top = Inches(0.3)
    
    mp1 = mtf.paragraphs[0]
    mp1.text = "📈 COMPILER BUILD METRICS"
    mp1.font.name = "Inter"
    mp1.font.size = Pt(15)
    mp1.font.bold = True
    mp1.font.color.rgb = EMERALD_GREEN
    
    mp2 = mtf.add_paragraph()
    mp2.text = "\n•  Turbopack Compile: 3.2 seconds\n•  TypeScript Validation: 123ms\n•  Runtime warnings: Zero (0)\n•  Runtime errors: Zero (0)\n\n•  User Interface Vetting:\nPremium dark-mode responsive glassmorphic cards with dynamic hover micro-animations fully compatible across mobile and widescreen viewports."
    mp2.font.name = "Inter"
    mp2.font.size = Pt(11)
    mp2.font.color.rgb = TEXT_MAIN
    mp2.space_before = Pt(5)

    # Save presentation
    save_path = "/Users/sridhargs/Documents/Antigravity/peer-bridge/PeerBridge_Technical_Masterclass_Deck.pptx"
    prs.save(save_path)
    print(f"Presentation created successfully at: {save_path}")

if __name__ == "__main__":
    create_presentation()
